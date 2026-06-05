import asyncio
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from fastapi import HTTPException, UploadFile, status

from schemas.rag import SegmentCreate
from settings import get_settings


@dataclass(frozen=True)
class DocumentTextExtractionResult:
    text: str
    segments: list[SegmentCreate]
    source_type: str


class DocumentTextExtractionService:
    """
    기능 요약: PDF/PPT/PPTX 업로드 파일에서 텍스트를 추출해 ingestion용 segment로 변환한다.

    기능 흐름:
        1. 업로드 파일을 임시 경로에 저장한다.
        2. 확장자별 loader로 페이지/슬라이드 텍스트를 추출한다.
        3. 각 페이지/슬라이드를 SegmentCreate로 변환하고 source range 컬럼 값을 채운다.

    파라미터 예시:
        file: FastAPI UploadFile("lecture.pdf")
        file_name: "lecture.pdf"
    """

    async def extract_upload(
        self,
        file: UploadFile,
        file_name: str,
    ) -> DocumentTextExtractionResult:
        suffix = Path(file_name).suffix.lower()
        if suffix not in {".pdf", ".ppt", ".pptx"}:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Unsupported document file type. Allowed extensions: .pdf, .ppt, .pptx",
            )

        with tempfile.TemporaryDirectory() as temp_dir_name:
            input_path = Path(temp_dir_name) / f"input{suffix}"
            await self._save_upload(file, input_path)

            if suffix == ".pdf":
                segments = await self._extract_pdf(input_path)
                source_type = "pdf"
            elif suffix == ".pptx":
                segments = await self._extract_pptx(input_path)
                source_type = "ppt"
            else:
                converted_path = await self._convert_legacy_ppt(input_path)
                segments = await self._extract_pptx(converted_path)
                source_type = "ppt"

        text = "\n".join(segment.text for segment in segments if segment.text.strip())
        if not text.strip():
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail="Document text extraction result is empty.",
            )

        return DocumentTextExtractionResult(
            text=text,
            segments=segments,
            source_type=source_type,
        )

    async def _save_upload(self, file: UploadFile, input_path: Path) -> None:
        with input_path.open("wb") as output_file:
            while chunk := await file.read(1024 * 1024):
                output_file.write(chunk)

        if input_path.stat().st_size == 0:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Uploaded document file is empty.",
            )

    async def _extract_pdf(self, path: Path) -> list[SegmentCreate]:
        try:
            from langchain_community.document_loaders import PyPDFLoader
        except Exception as exc:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="langchain-community and pypdf are required for PDF upload.",
            ) from exc

        loader = PyPDFLoader(str(path))
        documents = await asyncio.to_thread(loader.load)

        segments: list[SegmentCreate] = []
        for index, document in enumerate(documents):
            text = str(getattr(document, "page_content", "") or "").strip()
            if not text:
                continue
            metadata = getattr(document, "metadata", {}) or {}
            page_number = self._pdf_page_number(metadata, fallback=index + 1)
            segments.append(
                SegmentCreate(
                    segment_index=len(segments),
                    start_seconds=float(len(segments)),
                    end_seconds=float(len(segments) + 1),
                    text=text,
                    raw_metadata={
                        "loader": "PyPDFLoader",
                        "page": page_number,
                    },
                    source_type="pdf",
                    source_page_start=page_number,
                    source_page_end=page_number,
                )
            )

        return segments

    async def _extract_pptx(self, path: Path) -> list[SegmentCreate]:
        try:
            from pptx import Presentation
        except Exception as exc:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="python-pptx is required for PPTX upload.",
            ) from exc

        presentation = await asyncio.to_thread(Presentation, str(path))
        segments: list[SegmentCreate] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            text = self._slide_text(slide)
            if not text:
                continue
            segments.append(
                SegmentCreate(
                    segment_index=len(segments),
                    start_seconds=float(len(segments)),
                    end_seconds=float(len(segments) + 1),
                    text=text,
                    raw_metadata={
                        "loader": "python-pptx",
                        "slide": slide_index,
                    },
                    source_type="ppt",
                    source_slide_start=slide_index,
                    source_slide_end=slide_index,
                )
            )
        return segments

    async def _convert_legacy_ppt(self, path: Path) -> Path:
        settings = get_settings()
        if not settings.libreoffice_path:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail="Legacy PPT conversion is not configured.",
            )

        output_dir = path.parent / "converted"
        output_dir.mkdir(exist_ok=True)
        command = [
            settings.libreoffice_path,
            "--headless",
            "--convert-to",
            "pptx",
            "--outdir",
            str(output_dir),
            str(path),
        ]
        try:
            await asyncio.to_thread(
                subprocess.run,
                command,
                check=True,
                capture_output=True,
                text=True,
                timeout=settings.document_conversion_timeout_seconds,
            )
        except subprocess.SubprocessError as exc:
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail="Legacy PPT conversion failed.",
            ) from exc

        converted_path = output_dir / f"{path.stem}.pptx"
        if not converted_path.exists():
            raise HTTPException(
                status_code=status.HTTP_502_BAD_GATEWAY,
                detail="Legacy PPT conversion did not produce a PPTX file.",
            )
        return converted_path

    def _pdf_page_number(self, metadata: dict[str, Any], fallback: int) -> int:
        page = metadata.get("page")
        try:
            return int(page) + 1
        except (TypeError, ValueError):
            return fallback

    def _slide_text(self, slide: Any) -> str:
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape_text = str(getattr(shape, "text", "") or "").strip()
                if shape_text:
                    texts.append(shape_text)
            if getattr(shape, "has_table", False):
                texts.extend(self._table_texts(shape.table))
        return "\n".join(texts).strip()

    def _table_texts(self, table: Any) -> list[str]:
        texts: list[str] = []
        for row in table.rows:
            values = [
                str(cell.text or "").strip()
                for cell in row.cells
                if str(cell.text or "").strip()
            ]
            if values:
                texts.append(" | ".join(values))
        return texts
